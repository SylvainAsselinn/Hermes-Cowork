"""
Hermes Cowork - Main FastAPI Application
Full-featured backend for the Cowork dashboard
"""
from fastapi import FastAPI, WebSocket, WebSocketDisconnect, Depends, HTTPException, Query
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse, JSONResponse
from contextlib import asynccontextmanager
from datetime import datetime
from typing import List, Optional, Dict, Any
import asyncio
import json
import os
import sys

# Add parent to path for imports
sys.path.append(os.path.dirname(os.path.abspath(__file__)))

from database import Base, engine, get_db, Task, FileOperation, Agent, Setting, ChatMessage, TaskStatus, TaskPriority
from models import *
from file_manager import file_manager
from document_generator import generate_document, excel_generator, pptx_generator, pdf_generator
from subagent_manager import subagent_manager, SubAgent, SubTask, AgentStatus
from telegram_notifier import telegram_notifier

# Lifespan context manager
@asynccontextmanager
async def lifespan(app: FastAPI):
    # Startup
    print("🚀 Hermes Cowork starting...")
    # Create tables
    Base.metadata.create_all(bind=engine)
    # Initialize settings
    db = next(get_db())
    default_settings = {
        "max_agents": "8",
        "default_timeout": "300",
        "delete_protection": "true",
        "notifications_enabled": "true",
        "auto_backup": "true",
    }
    for key, value in default_settings.items():
        if not db.query(Setting).filter(Setting.key == key).first():
            db.add(Setting(key=key, value=value))
    db.commit()
    db.close()
    
    yield
    
    # Shutdown
    print("👋 Hermes Cowork shutting down...")


# Create app
app = FastAPI(
    title="Hermes Cowork API",
    description="Backend API for Hermes Cowork Dashboard",
    version="1.0.0",
    lifespan=lifespan,
)

# CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# WebSocket connection manager
class ConnectionManager:
    def __init__(self):
        self.active_connections: List[WebSocket] = []
    
    async def connect(self, websocket: WebSocket):
        await websocket.accept()
        self.active_connections.append(websocket)
    
    def disconnect(self, websocket: WebSocket):
        self.active_connections.remove(websocket)
    
    async def broadcast(self, message: dict):
        for connection in self.active_connections:
            try:
                await connection.send_json(message)
            except:
                pass

manager = ConnectionManager()


# ============ HEALTH & STATUS ============

@app.get("/")
async def root():
    """Root endpoint"""
    return {"name": "Hermes Cowork API", "version": "1.0.0", "status": "running"}


@app.get("/api/health")
async def health_check():
    """Health check endpoint"""
    return {
        "status": "healthy",
        "timestamp": datetime.utcnow().isoformat(),
        "services": {
            "database": "ok",
            "file_manager": "ok",
            "telegram": "ok" if telegram_notifier.token else "not_configured",
        }
    }


@app.get("/api/stats", response_model=DashboardStats)
async def get_dashboard_stats(db = Depends(get_db)):
    """Get dashboard statistics"""
    from datetime import timedelta
    
    tasks = db.query(Task).all()
    agents = subagent_manager.get_all_agents()
    
    # Calculate uptime (simplified)
    uptime = "0h 0m"
    if agents:
        oldest = min(a.started_at for a in agents)
        delta = datetime.utcnow() - oldest
        hours, remainder = divmod(int(delta.total_seconds()), 3600)
        minutes, _ = divmod(remainder, 60)
        uptime = f"{hours}h {minutes}m"
    
    return DashboardStats(
        tasks_total=len(tasks),
        tasks_pending=len([t for t in tasks if t.status == TaskStatus.PENDING]),
        tasks_in_progress=len([t for t in tasks if t.status == TaskStatus.IN_PROGRESS]),
        tasks_completed=len([t for t in tasks if t.status == TaskStatus.COMPLETED]),
        tasks_failed=len([t for t in tasks if t.status == TaskStatus.FAILED]),
        agents_active=len([a for a in agents if a.status == AgentStatus.ACTIVE]),
        files_processed_today=0,  # TODO: implement
        storage_used="N/A",  # TODO: implement
        uptime=uptime,
    )


# ============ TASK ENDPOINTS ============

@app.get("/api/tasks", response_model=List[TaskResponse])
async def list_tasks(
    status: Optional[TaskStatusEnum] = None,
    priority: Optional[TaskPriorityEnum] = None,
    limit: int = Query(50, ge=1, le=200),
    offset: int = Query(0, ge=0),
    db = Depends(get_db),
):
    """List all tasks with optional filtering"""
    query = db.query(Task)
    
    if status:
        query = query.filter(Task.status == status.value)
    if priority:
        query = query.filter(Task.priority == priority.value)
    
    tasks = query.order_by(Task.created_at.desc()).offset(offset).limit(limit).all()
    
    return [
        TaskResponse(
            id=t.id,
            title=t.title,
            description=t.description,
            status=t.status.value,
            priority=t.priority.value,
            progress=t.progress,
            created_at=t.created_at,
            started_at=t.started_at,
            completed_at=t.completed_at,
            estimated_duration=t.estimated_duration,
            actual_duration=t.actual_duration,
            parent_id=t.parent_id,
            result=t.result,
            error_message=t.error_message,
            input_files=json.loads(t.input_files) if t.input_files else [],
            output_files=json.loads(t.output_files) if t.output_files else [],
            agent_id=t.agent_id,
            subtasks=[],
        )
        for t in tasks
    ]


@app.post("/api/tasks", response_model=TaskResponse)
async def create_task(task: TaskCreate, db = Depends(get_db)):
    """Create a new task"""
    db_task = Task(
        title=task.title,
        description=task.description,
        status=TaskStatus.PENDING,
        priority=TaskPriority(task.priority.value),
        input_files=json.dumps(task.input_files),
        estimated_duration=task.estimated_duration,
        parent_id=task.parent_id,
    )
    db.add(db_task)
    db.commit()
    db.refresh(db_task)
    
    # Broadcast to WebSocket clients
    await manager.broadcast({
        "type": "task_created",
        "data": {"id": db_task.id, "title": db_task.title},
    })
    
    # Send Telegram notification
    await telegram_notifier.notify_task_started(db_task.id, db_task.title, task.estimated_duration)
    
    return TaskResponse(
        id=db_task.id,
        title=db_task.title,
        description=db_task.description,
        status=db_task.status.value,
        priority=db_task.priority.value,
        progress=db_task.progress,
        created_at=db_task.created_at,
        started_at=db_task.started_at,
        completed_at=db_task.completed_at,
        estimated_duration=db_task.estimated_duration,
        actual_duration=db_task.actual_duration,
        parent_id=db_task.parent_id,
        result=db_task.result,
        error_message=db_task.error_message,
        input_files=json.loads(db_task.input_files) if db_task.input_files else [],
        output_files=json.loads(db_task.output_files) if db_task.output_files else [],
        agent_id=db_task.agent_id,
        subtasks=[],
    )


@app.get("/api/tasks/{task_id}", response_model=TaskResponse)
async def get_task(task_id: int, db = Depends(get_db)):
    """Get a specific task"""
    task = db.query(Task).filter(Task.id == task_id).first()
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    
    # Get subtasks
    subtasks = db.query(Task).filter(Task.parent_id == task_id).all()
    
    return TaskResponse(
        id=task.id,
        title=task.title,
        description=task.description,
        status=task.status.value,
        priority=task.priority.value,
        progress=task.progress,
        created_at=task.created_at,
        started_at=task.started_at,
        completed_at=task.completed_at,
        estimated_duration=task.estimated_duration,
        actual_duration=task.actual_duration,
        parent_id=task.parent_id,
        result=task.result,
        error_message=task.error_message,
        input_files=json.loads(task.input_files) if task.input_files else [],
        output_files=json.loads(task.output_files) if task.output_files else [],
        agent_id=task.agent_id,
        subtasks=[
            TaskResponse(
                id=st.id,
                title=st.title,
                status=st.status.value,
                priority=st.priority.value,
                progress=st.progress,
                created_at=st.created_at,
                started_at=st.started_at,
                completed_at=st.completed_at,
                estimated_duration=st.estimated_duration,
                actual_duration=st.actual_duration,
                parent_id=st.parent_id,
                result=st.result,
                error_message=st.error_message,
                input_files=[],
                output_files=[],
                agent_id=st.agent_id,
                subtasks=[],
            )
            for st in subtasks
        ],
    )


@app.patch("/api/tasks/{task_id}", response_model=TaskResponse)
async def update_task(task_id: int, update: TaskUpdate, db = Depends(get_db)):
    """Update a task"""
    task = db.query(Task).filter(Task.id == task_id).first()
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    
    if update.title is not None:
        task.title = update.title
    if update.description is not None:
        task.description = update.description
    if update.status is not None:
        task.status = TaskStatus(update.status.value)
        if update.status == TaskStatusEnum.IN_PROGRESS and not task.started_at:
            task.started_at = datetime.utcnow()
        elif update.status in [TaskStatusEnum.COMPLETED, TaskStatusEnum.FAILED, TaskStatusEnum.CANCELLED]:
            task.completed_at = datetime.utcnow()
            if task.started_at:
                task.actual_duration = int((task.completed_at - task.started_at).total_seconds())
    if update.progress is not None:
        task.progress = update.progress
    if update.result is not None:
        task.result = update.result
    if update.error_message is not None:
        task.error_message = update.error_message
    
    db.commit()
    
    # Broadcast update
    await manager.broadcast({
        "type": "task_updated",
        "data": {"id": task.id, "status": task.status.value, "progress": task.progress},
    })
    
    return TaskResponse(
        id=task.id,
        title=task.title,
        description=task.description,
        status=task.status.value,
        priority=task.priority.value,
        progress=task.progress,
        created_at=task.created_at,
        started_at=task.started_at,
        completed_at=task.completed_at,
        estimated_duration=task.estimated_duration,
        actual_duration=task.actual_duration,
        parent_id=task.parent_id,
        result=task.result,
        error_message=task.error_message,
        input_files=json.loads(task.input_files) if task.input_files else [],
        output_files=json.loads(task.output_files) if task.output_files else [],
        agent_id=task.agent_id,
        subtasks=[],
    )


@app.delete("/api/tasks/{task_id}")
async def delete_task(task_id: int, db = Depends(get_db)):
    """Delete a task"""
    task = db.query(Task).filter(Task.id == task_id).first()
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    
    db.delete(task)
    db.commit()
    
    await manager.broadcast({
        "type": "task_deleted",
        "data": {"id": task_id},
    })
    
    return {"success": True, "message": "Task deleted"}


# ============ FILE ENDPOINTS ============

@app.get("/api/files")
async def list_files(path: str = None):
    """List files in a directory"""
    return file_manager.list_directory(path)


@app.get("/api/files/info")
async def get_file_info(path: str):
    """Get file information"""
    return file_manager.get_file_info(path)


@app.get("/api/files/read")
async def read_file_content(path: str, limit: int = 500):
    """Read file content"""
    return file_manager.read_file(path, limit)


@app.post("/api/files/write")
async def write_file_content(path: str, content: str, overwrite: bool = False):
    """Write content to a file"""
    return file_manager.write_file(path, content, overwrite)


@app.post("/api/files/delete/check")
async def check_delete_files(paths: List[str]):
    """Check what would be deleted (before confirmation)"""
    return file_manager.delete_check(paths)


@app.post("/api/files/delete")
async def delete_files(paths: List[str], backup: bool = True):
    """Delete files with optional backup"""
    result = file_manager.delete_files(paths, backup)
    
    # Broadcast update
    await manager.broadcast({
        "type": "files_deleted",
        "data": result,
    })
    
    return result


@app.post("/api/files/copy")
async def copy_file(source: str, destination: str):
    """Copy a file"""
    return file_manager.copy_file(source, destination)


@app.post("/api/files/move")
async def move_file(source: str, destination: str):
    """Move a file"""
    return file_manager.move_file(source, destination)


@app.get("/api/files/search")
async def search_files(pattern: str, directory: str = None, file_type: str = None):
    """Search for files"""
    return file_manager.search_files(pattern, directory, file_type)


@app.get("/api/files/quick-access")
async def get_quick_access():
    """Get quick access directories"""
    return file_manager.get_quick_access_dirs()


# ============ DOCUMENT GENERATION ENDPOINTS ============

@app.post("/api/documents/excel")
async def create_excel(title: str, data: Dict[str, Any], output_path: str = None):
    """Create an Excel document"""
    return excel_generator.create_from_data(title, data, output_path)


@app.post("/api/documents/powerpoint")
async def create_powerpoint(title: str, slides: List[Dict[str, Any]], output_path: str = None):
    """Create a PowerPoint presentation"""
    return pptx_generator.create_from_content(title, slides, output_path)


@app.post("/api/documents/pdf")
async def create_pdf(title: str, content: Dict[str, Any], output_path: str = None):
    """Create a PDF document"""
    return pdf_generator.create_report(title, content, output_path)


@app.get("/api/documents/download/{filename}")
async def download_document(filename: str):
    """Download a generated document"""
    file_path = os.path.join(os.path.expanduser("~/.hermes/cowork/data/documents"), filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="Document not found")
    return FileResponse(file_path, filename=filename)


# ============ SUB-AGENT ENDPOINTS ============

@app.get("/api/agents")
async def list_agents():
    """List all sub-agents"""
    return subagent_manager.get_status_summary()


@app.post("/api/agents/spawn")
async def spawn_agent(name: str = None, agent_prefix: str = "subagent"):
    """Spawn a new sub-agent"""
    try:
        agent = subagent_manager.spawn_agent(name, agent_prefix)
        await manager.broadcast({
            "type": "agent_spawned",
            "data": agent.to_dict(),
        })
        return agent.to_dict()
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


@app.get("/api/agents/{agent_id}")
async def get_agent(agent_id: str):
    """Get a specific agent"""
    agent = subagent_manager.get_agent(agent_id)
    if not agent:
        raise HTTPException(status_code=404, detail="Agent not found")
    return agent.to_dict()


@app.delete("/api/agents/{agent_id}")
async def kill_agent(agent_id: str):
    """Kill a sub-agent"""
    if subagent_manager.kill_agent(agent_id):
        await manager.broadcast({
            "type": "agent_killed",
            "data": {"agent_id": agent_id},
        })
        return {"success": True}
    raise HTTPException(status_code=404, detail="Agent not found")


@app.post("/api/tasks/{task_id}/spawn-agents")
async def spawn_agents_for_task(task_id: int, agent_count: int = 4):
    """Spawn agents for a task"""
    agents = subagent_manager.spawn_for_task(task_id, agent_count * 10, agent_count)
    await manager.broadcast({
        "type": "agents_spawned",
        "data": {"task_id": task_id, "agents": [a.to_dict() for a in agents]},
    })
    return [a.to_dict() for a in agents]


# ============ CHAT ENDPOINTS ============

@app.get("/api/chat/history")
async def get_chat_history(limit: int = 50, db = Depends(get_db)):
    """Get chat history"""
    messages = db.query(ChatMessage).order_by(ChatMessage.timestamp.desc()).limit(limit).all()
    return [
ChatMessageResponse(
 id=m.id,
 role=m.role,
 content=m.content,
 timestamp=m.timestamp,
 task_id=m.task_id,
 metadata=json.loads(m.extra_data) if m.extra_data else None,
 )
        for m in reversed(messages)
    ]


@app.post("/api/chat")
async def send_chat_message(message: ChatMessageCreate, db = Depends(get_db)):
    """Send a chat message"""
    db_message = ChatMessage(
        role=message.role,
        content=message.content,
        task_id=message.task_id,
        extra_data=json.dumps(message.metadata) if message.metadata else None,
    )
    db.add(db_message)
    db.commit()
    db.refresh(db_message)
    
    await manager.broadcast({
        "type": "chat_message",
        "data": {
            "id": db_message.id,
            "role": db_message.role,
            "content": db_message.content,
            "timestamp": db_message.timestamp.isoformat(),
        },
    })
    
    return ChatMessageResponse(
        id=db_message.id,
        role=db_message.role,
        content=db_message.content,
        timestamp=db_message.timestamp,
        task_id=db_message.task_id,
        metadata=json.loads(db_message.extra_data) if db_message.extra_data else None,
    )


# ============ SETTINGS ENDPOINTS ============

@app.get("/api/settings")
async def get_settings(db = Depends(get_db)):
    """Get all settings"""
    settings = db.query(Setting).all()
    return {s.key: s.value for s in settings}


@app.put("/api/settings")
async def update_setting(setting: SettingUpdate, db = Depends(get_db)):
    """Update a setting"""
    db_setting = db.query(Setting).filter(Setting.key == setting.key).first()
    if not db_setting:
        db_setting = Setting(key=setting.key, value=setting.value)
        db.add(db_setting)
    else:
        db_setting.value = setting.value
        db_setting.updated_at = datetime.utcnow()
    db.commit()
    return {"key": setting.key, "value": setting.value}


# ============ TELEGRAM NOTIFICATION ENDPOINTS ============

@app.post("/api/notify/telegram")
async def send_telegram_notification(notification: TelegramNotification):
    """Send a Telegram notification"""
    return await telegram_notifier.send_message(
        notification.message,
        buttons=notification.buttons,
    )


# ============ WEBSOCKET ============

@app.websocket("/ws")
async def websocket_endpoint(websocket: WebSocket):
    """WebSocket endpoint for real-time updates"""
    await manager.connect(websocket)
    try:
        while True:
            data = await websocket.receive_text()
            # Handle incoming WebSocket messages
            try:
                message = json.loads(data)
                # Echo back for now, could handle commands
                await websocket.send_json({
                    "type": "echo",
                    "data": message,
                    "timestamp": datetime.utcnow().isoformat(),
                })
            except json.JSONDecodeError:
                await websocket.send_json({
                    "type": "error",
                    "data": {"error": "Invalid JSON"},
                })
    except WebSocketDisconnect:
        manager.disconnect(websocket)



# ============ FILE PREVIEW ENDPOINT ============

def extract_docx_content(path: str) -> dict:
    """Extract text content from a .docx file"""
    try:
        from docx import Document
        doc = Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        tables = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables.append(table_data)
        return {
            "type": "docx",
            "paragraphs": paragraphs,
            "tables": tables,
            "total_paragraphs": len(paragraphs),
            "total_tables": len(tables)
        }
    except Exception as e:
        return {"error": str(e)}

def extract_xlsx_content(path: str) -> dict:
    """Extract content from a .xlsx file"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(path, read_only=True, data_only=True)
        sheets = {}
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []
            max_rows = min(sheet.max_row, 100)  # Limit to 100 rows
            max_cols = min(sheet.max_column, 20)  # Limit to 20 columns
            for row in sheet.iter_rows(min_row=1, max_row=max_rows, max_col=max_cols, values_only=True):
                rows.append([str(cell) if cell is not None else "" for cell in row])
            sheets[sheet_name] = rows
        wb.close()
        return {
            "type": "xlsx",
            "sheets": sheets,
            "sheet_names": list(sheets.keys())
        }
    except Exception as e:
        return {"error": str(e)}

def extract_pptx_content(path: str) -> dict:
    """Extract content from a .pptx file"""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_content = {"slide_number": i + 1, "shapes": []}
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content["shapes"].append({
                        "type": "text",
                        "content": shape.text
                    })
                if hasattr(shape, "image"):
                    slide_content["shapes"].append({
                        "type": "image",
                        "content_type": shape.image.content_type
                    })
            slides.append(slide_content)
        return {
            "type": "pptx",
            "slides": slides,
            "total_slides": len(slides)
        }
    except Exception as e:
        return {"error": str(e)}

@app.get("/api/files/preview")
async def preview_file(path: str):
    """Serve a file for preview (images, PDFs, Office files, etc.)"""
    import mimetypes
    mimetypes.init()
    
    if not os.path.exists(path):
        raise HTTPException(status_code=404, detail="File not found")
    
    # Check if it's a file (not a directory)
    if not os.path.isfile(path):
        raise HTTPException(status_code=400, detail="Not a file")
    
    # Get mime type
    mime_type, _ = mimetypes.guess_type(path)
    if mime_type is None:
        mime_type = "application/octet-stream"
    
    # For text files, return as plain text with proper encoding
    text_extensions = ['.txt', '.md', '.py', '.js', '.ts', '.jsx', '.tsx', '.json', 
    '.yaml', '.yml', '.xml', '.html', '.css', '.scss', '.sass', '.less',
    '.sh', '.bash', '.zsh', '.fish', '.ps1', '.bat', '.cmd',
    '.java', '.cpp', '.c', '.h', '.hpp', '.cs', '.go', '.rs', '.rb',
    '.php', '.sql', '.swift', '.kt', '.scala', '.lua', '.r', '.m',
    '.toml', '.ini', '.cfg', '.conf', '.log', '.env', '.gitignore',
    '.dockerfile', '.makefile', '.cmake']
    
    _, ext = os.path.splitext(path.lower())
    
    # Office files
    if ext == '.docx':
        result = extract_docx_content(path)
        return JSONResponse(result)
    
    if ext in ['.xlsx', '.xls']:
        result = extract_xlsx_content(path)
        return JSONResponse(result)
    
    if ext in ['.pptx', '.ppt']:
        result = extract_pptx_content(path)
        return JSONResponse(result)
    
    if ext in text_extensions or mime_type.startswith('text/'):
        try:
            with open(path, 'r', encoding='utf-8', errors='replace') as f:
                file_content = f.read(50000)  # Limit to 50KB
            return JSONResponse({
                "type": "text",
                "content": file_content,
                "extension": ext,
                "path": path
            })
        except Exception as e:
            return JSONResponse({"error": str(e)}, status_code=500)
    
    # For binary files (images, PDFs), return the file
    return FileResponse(
        path,
        media_type=mime_type,
        filename=os.path.basename(path)
    )

# ============ STATIC FILES (for frontend) ============

# Mount static files for frontend (will be created in next step)
FRONTEND_DIST = os.path.expanduser("~/.hermes/cowork/frontend/dist")
if os.path.exists(FRONTEND_DIST):
    app.mount("/assets", StaticFiles(directory=os.path.join(FRONTEND_DIST, "assets")), name="assets")
    
    @app.get("/{path:path}")
    async def serve_frontend(path: str):
        """Serve frontend for all other routes"""
        file_path = os.path.join(FRONTEND_DIST, path)
        if os.path.exists(file_path) and os.path.isfile(file_path):
            return FileResponse(file_path)
        return FileResponse(os.path.join(FRONTEND_DIST, "index.html"))


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8000)
