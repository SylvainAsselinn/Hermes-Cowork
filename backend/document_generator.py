"""
Hermes Cowork - Document Generator
Creates Excel, PowerPoint, PDF, and Word documents
"""
import os
from datetime import datetime
from typing import Dict, Any, List, Optional
import json

# Excel
from openpyxl import Workbook
from openpyxl.styles import Font, Fill, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, PieChart, LineChart, Reference
from openpyxl.utils import get_column_letter

# PowerPoint
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# PDF
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch, cm
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak, Image
from reportlab.lib.enums import TA_CENTER, TA_LEFT

# Output directory
OUTPUT_DIR = os.path.expanduser("~/.hermes/cowork/data/documents")
os.makedirs(OUTPUT_DIR, exist_ok=True)


class ExcelGenerator:
    """Generate Excel files with formatting and charts"""
    
    def __init__(self):
        self.wb = Workbook()
        self.ws = self.wb.active
    
    def create_from_data(self, title: str, data: Dict[str, Any], output_path: str = None) -> Dict[str, Any]:
        """Create Excel from structured data"""
        try:
            self.ws.title = title[:31]  # Excel sheet name limit
            
            # Styles
            header_font = Font(bold=True, color="FFFFFF")
            header_fill = PatternFill(start_color="2E7D32", end_color="2E7D32", fill_type="solid")
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # Check data format
            if "headers" in data and "rows" in data:
                # Table format
                headers = data["headers"]
                rows = data["rows"]
                
                # Write headers
                for col, header in enumerate(headers, 1):
                    cell = self.ws.cell(row=1, column=col, value=header)
                    cell.font = header_font
                    cell.fill = header_fill
                    cell.alignment = Alignment(horizontal='center')
                    cell.border = border
                
                # Write data
                for row_idx, row_data in enumerate(rows, 2):
                    for col_idx, value in enumerate(row_data, 1):
                        cell = self.ws.cell(row=row_idx, column=col_idx, value=value)
                        cell.border = border
                        if isinstance(value, (int, float)) and col_idx > 1:
                            cell.alignment = Alignment(horizontal='right')
                
                # Auto-adjust column widths
                for col in range(1, len(headers) + 1):
                    max_length = max(
                        len(str(headers[col-1])),
                        max(len(str(row[col-1])) for row in rows) if rows else 0
                    )
                    self.ws.column_dimensions[get_column_letter(col)].width = min(max_length + 2, 50)
                
                # Add chart if numeric data
                if len(rows) > 1 and all(isinstance(r[1], (int, float)) for r in rows[:5]):
                    chart = BarChart()
                    chart.type = "col"
                    chart.style = 10
                    chart.title = title
                    chart.y_axis.title = headers[1] if len(headers) > 1 else "Valeur"
                    chart.x_axis.title = headers[0]
                    
                    data_ref = Reference(self.ws, min_col=2, min_row=1, max_row=len(rows)+1)
                    cats = Reference(self.ws, min_col=1, min_row=2, max_row=len(rows)+1)
                    chart.add_data(data_ref, titles_from_data=True)
                    chart.set_categories(cats)
                    chart.shape = 4
                    self.ws.add_chart(chart, f"{get_column_letter(len(headers)+2)}1")
            
            elif "data" in data:
                # Key-value format
                row = 1
                for key, value in data["data"].items():
                    self.ws.cell(row=row, column=1, value=key).font = Font(bold=True)
                    self.ws.cell(row=row, column=2, value=value)
                    row += 1
                
                self.ws.column_dimensions['A'].width = 30
                self.ws.column_dimensions['B'].width = 50
            
            # Save
            if output_path is None:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_path = os.path.join(OUTPUT_DIR, f"{title.replace(' ', '_')}_{timestamp}.xlsx")
            
            self.wb.save(output_path)
            
            return {
                "success": True,
                "file_path": output_path,
                "file_name": os.path.basename(output_path),
                "file_size": os.path.getsize(output_path),
            }
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def create_from_list(self, title: str, items: List[Dict[str, Any]], output_path: str = None) -> Dict[str, Any]:
        """Create Excel from list of dictionaries"""
        if not items:
            return {"success": False, "error": "No data provided"}
        
        headers = list(items[0].keys())
        rows = [[item.get(h, "") for h in headers] for item in items]
        
        return self.create_from_data(title, {"headers": headers, "rows": rows}, output_path)


class PowerPointGenerator:
    """Generate PowerPoint presentations"""
    
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
    
    def create_from_content(self, title: str, slides: List[Dict[str, Any]], output_path: str = None) -> Dict[str, Any]:
        """Create PowerPoint from structured content"""
        try:
            # Title slide
            title_slide_layout = self.prs.slide_layouts[6]  # Blank
            slide = self.prs.slides.add_slide(title_slide_layout)
            
            # Title
            left = Inches(0.5)
            top = Inches(2.5)
            width = Inches(12.333)
            height = Inches(2)
            
            title_box = slide.shapes.add_textbox(left, top, width, height)
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = title
            title_para.font.size = Pt(44)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(46, 125, 50)
            title_para.alignment = PP_ALIGN.CENTER
            
            # Date
            date_box = slide.shapes.add_textbox(left, Inches(5), width, Inches(0.5))
            date_frame = date_box.text_frame
            date_para = date_frame.paragraphs[0]
            date_para.text = datetime.now().strftime("%d/%m/%Y")
            date_para.font.size = Pt(18)
            date_para.alignment = PP_ALIGN.CENTER
            
            # Content slides
            for slide_data in slides:
                self._add_content_slide(slide_data)
            
            # Save
            if output_path is None:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_path = os.path.join(OUTPUT_DIR, f"{title.replace(' ', '_')}_{timestamp}.pptx")
            
            self.prs.save(output_path)
            
            return {
                "success": True,
                "file_path": output_path,
                "file_name": os.path.basename(output_path),
                "file_size": os.path.getsize(output_path),
                "slides_count": len(self.prs.slides),
            }
        except Exception as e:
            return {"success": False, "error": str(e)}
    
    def _add_content_slide(self, slide_data: Dict[str, Any]):
        """Add a content slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        if "title" in slide_data:
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
            tf = title_box.text_frame
            p = tf.paragraphs[0]
            p.text = slide_data["title"]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = RGBColor(46, 125, 50)
        
        # Content
        if "content" in slide_data:
            content = slide_data["content"]
            
            if isinstance(content, list):
                # Bullet points
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
                tf = content_box.text_frame
                tf.word_wrap = True
                
                for i, item in enumerate(content):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
                    p.text = f"• {item}"
                    p.font.size = Pt(20)
                    p.space_after = Pt(12)
            
            elif isinstance(content, str):
                # Text content
                content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
                tf = content_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = content
                p.font.size = Pt(18)
        
        # Table
        if "table" in slide_data:
            table_data = slide_data["table"]
            rows = len(table_data)
            cols = len(table_data[0]) if table_data else 0
            
            if rows > 0 and cols > 0:
                x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(12), Inches(0.5 * rows)
                shape = slide.shapes.add_table(rows, cols, x, y, cx, cy)
                table = shape.table
                
                for i, row_data in enumerate(table_data):
                    for j, cell_value in enumerate(row_data):
                        cell = table.cell(i, j)
                        cell.text = str(cell_value)
                        if i == 0:  # Header
                            cell.text_frame.paragraphs[0].font.bold = True
    
    def create_simple(self, title: str, content: List[str], output_path: str = None) -> Dict[str, Any]:
        """Create simple presentation with title and bullet slides"""
        slides = [{"title": "Contenu", "content": content}]
        return self.create_from_content(title, slides, output_path)


class PDFGenerator:
    """Generate PDF documents"""
    
    def create_report(self, title: str, content: Dict[str, Any], output_path: str = None) -> Dict[str, Any]:
        """Create PDF report"""
        try:
            if output_path is None:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                output_path = os.path.join(OUTPUT_DIR, f"{title.replace(' ', '_')}_{timestamp}.pdf")
            
            doc = SimpleDocTemplate(output_path, pagesize=A4, 
                                   rightMargin=72, leftMargin=72,
                                   topMargin=72, bottomMargin=72)
            
            styles = getSampleStyleSheet()
            styles.add(ParagraphStyle(
                name='CustomTitle',
                parent=styles['Heading1'],
                fontSize=24,
                spaceAfter=30,
                alignment=TA_CENTER,
                textColor=colors.HexColor('#2E7D32')
            ))
            
            story = []
            
            # Title
            story.append(Paragraph(title, styles['CustomTitle']))
            story.append(Spacer(1, 12))
            
            # Date
            story.append(Paragraph(f"Date: {datetime.now().strftime('%d/%m/%Y')}", styles['Normal']))
            story.append(Spacer(1, 30))
            
            # Content
            if "sections" in content:
                for section in content["sections"]:
                    if "title" in section:
                        story.append(Paragraph(section["title"], styles['Heading2']))
                        story.append(Spacer(1, 12))
                    
                    if "text" in section:
                        story.append(Paragraph(section["text"], styles['Normal']))
                        story.append(Spacer(1, 12))
                    
                    if "items" in section:
                        for item in section["items"]:
                            story.append(Paragraph(f"• {item}", styles['Normal']))
                        story.append(Spacer(1, 12))
                    
                    if "table" in section:
                        table_data = section["table"]
                        if table_data:
                            t = Table(table_data)
                            t.setStyle(TableStyle([
                                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2E7D32')),
                                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                                ('FONTSIZE', (0, 0), (-1, 0), 12),
                                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                                ('TEXTCOLOR', (0, 1), (-1, -1), colors.black),
                                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                            ]))
                            story.append(t)
                            story.append(Spacer(1, 12))
            
            doc.build(story)
            
            return {
                "success": True,
                "file_path": output_path,
                "file_name": os.path.basename(output_path),
                "file_size": os.path.getsize(output_path),
            }
        except Exception as e:
            return {"success": False, "error": str(e)}


# Singleton instances
excel_generator = ExcelGenerator()
pptx_generator = PowerPointGenerator()
pdf_generator = PDFGenerator()


def generate_document(doc_type: str, title: str, data: Dict[str, Any], output_path: str = None) -> Dict[str, Any]:
    """Generate document based on type"""
    if doc_type.lower() in ["excel", "xlsx", "xls"]:
        return excel_generator.create_from_data(title, data, output_path)
    elif doc_type.lower() in ["powerpoint", "pptx", "ppt"]:
        slides = data.get("slides", [{"title": "Contenu", "content": data.get("items", [])}])
        return pptx_generator.create_from_content(title, slides, output_path)
    elif doc_type.lower() in ["pdf"]:
        return pdf_generator.create_report(title, data, output_path)
    else:
        return {"success": False, "error": f"Unknown document type: {doc_type}"}
